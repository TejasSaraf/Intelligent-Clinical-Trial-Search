#!/usr/bin/env python3
"""
Generate a PowerPoint presentation for Clinical Trial Search.
Run: pip install python-pptx && python scripts/create_presentation.py
"""
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    raise

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "ClinicalTrialSearch_Presentation.pptx"

# Theme colors
TITLE_COLOR = RgbColor(0x22, 0xC5, 0x5E)  # Green accent
DARK_BG = RgbColor(0x17, 0x24, 0x2E)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left, top, w, h = Inches(0.5), Inches(2), Inches(9), Inches(1.2)
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    if subtitle:
        tx2 = slide.shapes.add_textbox(left, Inches(3.2), w, Inches(1.5))
        tx2.text_frame.text = subtitle
        tx2.text_frame.paragraphs[0].font.size = Pt(20)
        tx2.text_frame.paragraphs[0].font.color.rgb = RgbColor(0x6B, 0x72, 0x80)
    return slide


def add_content_slide(prs, title, bullets, subtext=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left, top = Inches(0.5), Inches(0.6)
    tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    body = slide.shapes.add_textbox(left, Inches(1.5), Inches(9), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    if subtext:
        st = slide.shapes.add_textbox(left, Inches(6.2), Inches(9), Inches(0.6))
        st.text_frame.text = subtext
        st.text_frame.paragraphs[0].font.size = Pt(12)
        st.text_frame.paragraphs[0].font.color.rgb = RgbColor(0x9C, 0xA3, 0xAF)
    return slide


def add_two_column_slide(prs, title, left_col, right_col):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
    tx.text_frame.paragraphs[0].text = title
    tx.text_frame.paragraphs[0].font.size = Pt(28)
    tx.text_frame.paragraphs[0].font.bold = True
    tx.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    lb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5))
    lb.text_frame.word_wrap = True
    for i, item in enumerate(left_col):
        p = lb.text_frame.paragraphs[i] if i < len(lb.text_frame.paragraphs) else lb.text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    rb = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.3), Inches(5))
    rb.text_frame.word_wrap = True
    for i, item in enumerate(right_col):
        p = rb.text_frame.paragraphs[i] if i < len(rb.text_frame.paragraphs) else rb.text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Clinical Trial Search",
        "Discover clinical trials that match—not the other way around.",
    )

    add_content_slide(
        prs,
        "Problem & Solution",
        [
            "Finding relevant clinical trials is hard (keywords, filters, jargon)",
            "Users think in natural language, not structured filters",
            "Solution: Natural-language search with AI-powered summaries",
        ],
    )

    add_content_slide(
        prs,
        "Key Features",
        [
            "Natural language search (e.g. 'Phase 2 breast cancer trials')",
            "Spell correction for typos and medical terms",
            "AI summaries of results (RAG)",
            "Thinking message while search runs",
            "Show more pagination",
            "Full trial details modal",
        ],
    )

    add_two_column_slide(
        prs,
        "Tech Stack",
        ["Backend: Python, FastAPI", "Search: Elasticsearch", "AI: OpenAI (optional)", "RAG for summaries & spell correction"],
        ["Frontend: React 19, TypeScript", "Build: Vite", "Styling: Tailwind CSS", "API client for search, trial details"],
    )

    add_content_slide(
        prs,
        "Architecture",
        [
            "User query → Spell correction (LLM) → Query parser → Elasticsearch",
            "ES results → Summary generation (RAG) → Response",
            "Trial details: GET /trial/{nct_id} on demand",
        ],
        "Minimal list payload; full details fetched only when 'View details' is clicked.",
    )

    add_content_slide(
        prs,
        "API Endpoints",
        [
            "GET /search/{query} — Search trials (?page, ?size)",
            "GET /search/thinking/{query} — LLM 'thinking' message",
            "GET /trial/{nct_id} — Full trial details by NCT ID",
        ],
    )

    add_content_slide(
        prs,
        "UX Highlights",
        [
            "Chat-like interface: user query + assistant summary + results",
            "Clear progress: 'Thinking…' while search runs",
            "Load more with skeleton loading",
            "Modal for trial details (gray overlay, response-style background)",
        ],
    )

    add_content_slide(
        prs,
        "Setup & Run",
        [
            "Elasticsearch 8.x (Docker or local)",
            "Backend: Python venv, pip install, ingest data, uvicorn",
            "Frontend: npm install, npm run dev",
            "Optional: OPENAI_API_KEY for AI features",
        ],
    )

    add_title_slide(
        prs,
        "Thank You",
        "Questions?",
    )

    prs.save(OUTPUT)
    print(f"Presentation saved to {OUTPUT}")


if __name__ == "__main__":
    main()
